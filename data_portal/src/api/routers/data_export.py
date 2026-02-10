"""
바-1: 데이터 반출 API — CSV/JSON 다운로드 + IRB 승인 워크플로우
RFP 요구: 표준 인터페이스 기반 데이터 서비스 연동
"""
import os
import io
import csv
import json
import time
from datetime import datetime
from typing import Optional, List

from fastapi import APIRouter, HTTPException, Query, Depends
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field

router = APIRouter(prefix="/export", tags=["DataExport"])

# ── DB 연결 ──
async def _get_conn():
    from services.db_pool import get_pool
    pool = await get_pool()
    return await pool.acquire()

async def _rel(conn):
    from services.db_pool import get_pool
    pool = await get_pool()
    await pool.release(conn)

# ── Pydantic 모델 ──
class ExportRequest(BaseModel):
    sql: str = Field(..., min_length=5, max_length=5000)
    format: str = Field(default="csv", pattern=r"^(csv|json|xlsx|pptx|pdf)$")
    filename: Optional[str] = Field(None, max_length=200)
    purpose: str = Field(..., min_length=2, max_length=500)
    irb_number: Optional[str] = Field(None, max_length=50)

class ExportApprovalRequest(BaseModel):
    status: str = Field(..., pattern=r"^(approved|rejected)$")
    reviewer: str = Field(default="admin", max_length=50)
    comment: Optional[str] = Field(None, max_length=1000)

# ── 테이블 초기화 ──
_tbl_ok = False
async def _ensure_tables(conn):
    global _tbl_ok
    if _tbl_ok:
        return
    await conn.execute("""
        CREATE TABLE IF NOT EXISTS export_request (
            request_id SERIAL PRIMARY KEY,
            requester VARCHAR(50) NOT NULL DEFAULT 'demo',
            sql_query TEXT NOT NULL,
            format VARCHAR(10) DEFAULT 'csv',
            filename VARCHAR(200),
            purpose TEXT NOT NULL,
            irb_number VARCHAR(50),
            status VARCHAR(20) DEFAULT 'pending',
            row_count INTEGER,
            reviewer VARCHAR(50),
            review_comment TEXT,
            reviewed_at TIMESTAMPTZ,
            created_at TIMESTAMPTZ DEFAULT NOW()
        );
        CREATE INDEX IF NOT EXISTS idx_export_req_status ON export_request(status);
    """)
    _tbl_ok = True

MAX_EXPORT_ROWS = int(os.getenv("EXPORT_MAX_ROWS", "100000"))

# ── SQL 검증 (읽기 전용) ──
import re
_FORBIDDEN = ["INSERT","UPDATE","DELETE","DROP","CREATE","ALTER","TRUNCATE",
              "GRANT","REVOKE","EXECUTE","COPY","LOAD"]

def _validate_export_sql(sql: str) -> tuple:
    s = sql.upper().strip()
    if not (s.startswith("SELECT") or s.startswith("WITH")):
        return False, "SELECT/WITH 쿼리만 허용됩니다"
    for kw in _FORBIDDEN:
        if re.search(rf'\b{kw}\b', s):
            return False, f"금지 키워드: {kw}"
    if "--" in sql or "/*" in sql:
        return False, "SQL 주석 금지"
    return True, ""


# ══════════════════════════════════════════
# 포맷별 변환 헬퍼
# ══════════════════════════════════════════

def _export_xlsx(rows, columns, fname):
    """Excel (xlsx) 변환"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()
    ws = wb.active
    ws.title = "데이터"

    # 헤더 스타일
    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="2B579A", end_color="2B579A", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )

    for ci, col in enumerate(columns, 1):
        cell = ws.cell(row=1, column=ci, value=col)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")
        cell.border = thin_border

    for ri, row in enumerate(rows, 2):
        for ci, col in enumerate(columns, 1):
            val = row[col]
            if hasattr(val, "isoformat"):
                val = val.isoformat()
            cell = ws.cell(row=ri, column=ci, value=str(val) if val is not None else "")
            cell.border = thin_border

    # 열 너비 자동 조정
    for ci, col in enumerate(columns, 1):
        max_len = max(len(str(col)), *(len(str(row[col] or "")[:50]) for row in rows[:100]))
        ws.column_dimensions[ws.cell(row=1, column=ci).column_letter].width = min(max_len + 4, 50)

    # 필터 설정
    ws.auto_filter.ref = ws.dimensions

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{fname}.xlsx"'},
    )


def _export_pptx(rows, columns, fname):
    """PowerPoint (pptx) 변환 — 데이터 테이블 슬라이드"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 제목 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"서울아산병원 IDP — {fname}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2B, 0x57, 0x9A)

    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"총 {len(rows):,}건 · {len(columns)}컬럼 · {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # 데이터 슬라이드 (50행씩 분할)
    PAGE_SIZE = 50
    max_cols = min(len(columns), 12)  # PPT 폭 제한
    display_cols = columns[:max_cols]

    for page_start in range(0, len(rows), PAGE_SIZE):
        page_rows = rows[page_start:page_start + PAGE_SIZE]
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 페이지 번호
        page_num = page_start // PAGE_SIZE + 1
        total_pages = (len(rows) + PAGE_SIZE - 1) // PAGE_SIZE
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Page {page_num}/{total_pages} (행 {page_start+1}~{page_start+len(page_rows)})"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        # 테이블
        tbl_rows = len(page_rows) + 1
        tbl_cols = len(display_cols)
        left = Inches(0.3)
        top = Inches(0.7)
        width = Inches(12.7)
        height = Inches(6.5)

        table = slide.shapes.add_table(tbl_rows, tbl_cols, left, top, width, height).table

        # 헤더
        for ci, col in enumerate(display_cols):
            cell = table.cell(0, ci)
            cell.text = str(col)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2B, 0x57, 0x9A)

        # 데이터
        for ri, row in enumerate(page_rows, 1):
            for ci, col in enumerate(display_cols):
                cell = table.cell(ri, ci)
                val = row[col]
                if hasattr(val, "isoformat"):
                    val = val.isoformat()
                cell.text = str(val)[:100] if val is not None else ""
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(8)
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{fname}.pptx"'},
    )


def _export_pdf(rows, columns, fname):
    """PDF 변환 — 간단한 HTML→PDF (reportlab 없이 텍스트 기반)"""
    # HTML 테이블 생성 후 브라우저에서 인쇄 가능한 형태로 제공
    html_parts = [
        "<!DOCTYPE html><html><head><meta charset='utf-8'>",
        "<style>",
        "body{font-family:'Malgun Gothic',sans-serif;margin:20px}",
        "h1{color:#2B579A;font-size:18px}",
        "table{border-collapse:collapse;width:100%;font-size:11px}",
        "th{background:#2B579A;color:#fff;padding:6px 8px;text-align:left;border:1px solid #ccc}",
        "td{padding:4px 8px;border:1px solid #ddd}",
        "tr:nth-child(even){background:#f9f9f9}",
        ".meta{color:#666;font-size:12px;margin-bottom:10px}",
        "@media print{body{margin:0}h1{font-size:14px}}",
        "</style></head><body>",
        f"<h1>서울아산병원 IDP — {fname}</h1>",
        f"<p class='meta'>총 {len(rows):,}건 · {len(columns)}컬럼 · 생성일: {datetime.now().strftime('%Y-%m-%d %H:%M')}</p>",
        "<table><thead><tr>",
    ]
    for col in columns:
        html_parts.append(f"<th>{col}</th>")
    html_parts.append("</tr></thead><tbody>")

    for row in rows[:5000]:  # PDF는 5000행 제한
        html_parts.append("<tr>")
        for col in columns:
            val = row[col]
            if hasattr(val, "isoformat"):
                val = val.isoformat()
            html_parts.append(f"<td>{str(val) if val is not None else ''}</td>")
        html_parts.append("</tr>")
    html_parts.append("</tbody></table>")
    if len(rows) > 5000:
        html_parts.append(f"<p class='meta'>※ PDF는 상위 5,000건만 포함 (전체: {len(rows):,}건)</p>")
    html_parts.append("</body></html>")

    html = "".join(html_parts)
    return StreamingResponse(
        io.BytesIO(html.encode("utf-8")),
        media_type="text/html",
        headers={"Content-Disposition": f'attachment; filename="{fname}.html"'},
    )


# ══════════════════════════════════════════
# 엔드포인트
# ══════════════════════════════════════════

@router.post("/request")
async def create_export_request(body: ExportRequest):
    """반출 요청 생성 (IRB 번호 필수 시 검증)"""
    ok, msg = _validate_export_sql(body.sql)
    if not ok:
        raise HTTPException(400, msg)
    conn = await _get_conn()
    try:
        await _ensure_tables(conn)
        # IRB 필수 여부 확인
        irb_required = await conn.fetchval(
            "SELECT setting_value::text FROM po_system_setting WHERE setting_key = 'export.require_approval'"
        )
        if irb_required and irb_required.strip('"') == "true" and not body.irb_number:
            raise HTTPException(400, "데이터 반출에 IRB 승인 번호가 필요합니다")

        # 행 수 사전 추정
        count_sql = f"SELECT COUNT(*) FROM ({body.sql}) t"
        try:
            row_est = await conn.fetchval(count_sql)
        except Exception:
            row_est = None
        if row_est and row_est > MAX_EXPORT_ROWS:
            raise HTTPException(400, f"반출 최대 행 수 초과: {row_est:,} > {MAX_EXPORT_ROWS:,}")

        auto_approve = (irb_required and irb_required.strip('"') != "true")
        status = "approved" if auto_approve else "pending"

        rid = await conn.fetchval("""
            INSERT INTO export_request (requester, sql_query, format, filename, purpose, irb_number, status, row_count)
            VALUES ($1,$2,$3,$4,$5,$6,$7,$8) RETURNING request_id
        """, "demo", body.sql, body.format, body.filename or f"export_{int(time.time())}",
            body.purpose, body.irb_number, status, row_est)
        return {"request_id": rid, "status": status, "estimated_rows": row_est}
    finally:
        await _rel(conn)


@router.get("/requests")
async def list_export_requests(status: Optional[str] = None, limit: int = Query(50, le=200)):
    """반출 요청 목록"""
    conn = await _get_conn()
    try:
        await _ensure_tables(conn)
        q = "SELECT * FROM export_request"
        params = []
        if status:
            q += " WHERE status = $1"
            params.append(status)
        q += " ORDER BY created_at DESC LIMIT " + str(limit)
        rows = await conn.fetch(q, *params)
        return [dict(r) for r in rows]
    finally:
        await _rel(conn)


@router.put("/requests/{request_id}/review")
async def review_export_request(request_id: int, body: ExportApprovalRequest):
    """반출 요청 승인/거절"""
    conn = await _get_conn()
    try:
        await _ensure_tables(conn)
        existing = await conn.fetchrow("SELECT * FROM export_request WHERE request_id = $1", request_id)
        if not existing:
            raise HTTPException(404, "요청을 찾을 수 없습니다")
        if existing["status"] != "pending":
            raise HTTPException(400, f"이미 처리된 요청: {existing['status']}")
        await conn.execute("""
            UPDATE export_request SET status=$1, reviewer=$2, review_comment=$3, reviewed_at=NOW()
            WHERE request_id=$4
        """, body.status, body.reviewer, body.comment, request_id)
        return {"request_id": request_id, "status": body.status}
    finally:
        await _rel(conn)


@router.get("/download/{request_id}")
async def download_export(request_id: int):
    """승인된 반출 요청 데이터 다운로드 (CSV/JSON)"""
    conn = await _get_conn()
    try:
        await _ensure_tables(conn)
        req = await conn.fetchrow("SELECT * FROM export_request WHERE request_id = $1", request_id)
        if not req:
            raise HTTPException(404, "요청을 찾을 수 없습니다")
        if req["status"] != "approved":
            raise HTTPException(403, f"승인되지 않은 요청: {req['status']}")

        # SQL 실행
        sql = req["sql_query"]
        if "LIMIT" not in sql.upper():
            sql += f" LIMIT {MAX_EXPORT_ROWS}"

        rows = await conn.fetch(sql)
        if not rows:
            raise HTTPException(404, "결과가 없습니다")

        columns = list(rows[0].keys())
        fname = req["filename"] or f"export_{request_id}"

        fmt = req["format"]

        if fmt == "json":
            data = json.dumps([dict(r) for r in rows], ensure_ascii=False, default=str, indent=2)
            return StreamingResponse(
                io.BytesIO(data.encode("utf-8")),
                media_type="application/json",
                headers={"Content-Disposition": f'attachment; filename="{fname}.json"'}
            )
        elif fmt == "xlsx":
            return _export_xlsx(rows, columns, fname)
        elif fmt == "pptx":
            return _export_pptx(rows, columns, fname)
        elif fmt == "pdf":
            return _export_pdf(rows, columns, fname)
        else:
            buf = io.StringIO()
            writer = csv.writer(buf)
            writer.writerow(columns)
            for r in rows:
                writer.writerow([r[c] for c in columns])
            return StreamingResponse(
                io.BytesIO(buf.getvalue().encode("utf-8-sig")),
                media_type="text/csv",
                headers={"Content-Disposition": f'attachment; filename="{fname}.csv"'}
            )
    finally:
        await _rel(conn)


@router.get("/stats")
async def export_stats():
    """반출 통계"""
    conn = await _get_conn()
    try:
        await _ensure_tables(conn)
        stats = await conn.fetchrow("""
            SELECT
                COUNT(*) as total_requests,
                COUNT(*) FILTER (WHERE status='approved') as approved,
                COUNT(*) FILTER (WHERE status='pending') as pending,
                COUNT(*) FILTER (WHERE status='rejected') as rejected,
                COALESCE(SUM(row_count) FILTER (WHERE status='approved'), 0) as total_rows_exported
            FROM export_request
        """)
        return dict(stats)
    finally:
        await _rel(conn)
