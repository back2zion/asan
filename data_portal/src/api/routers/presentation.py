"""
프레젠테이션 자동 생성 API
- Paper2Slides (HKUDS) 연동
- GPU 서버의 Qwen3-32B 활용
- API: POST /generate → GET /status/{job_id} → GET /download/{job_id}
- PPTX 변환: 생성된 PDF를 이미지로 변환 후 python-pptx로 PPTX 생성
"""
import io
import httpx
from typing import Optional, List
from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Emu

router = APIRouter(prefix="/presentation", tags=["Presentation"])

# Paper2Slides API (SSH 터널: localhost:29001 -> GPU서버:9001)
PAPER2SLIDES_API_URL = "http://localhost:29001"


class JobResponse(BaseModel):
    job_id: str
    status: str
    message: Optional[str] = None


class JobStatus(BaseModel):
    job_id: str
    status: str
    output_type: str
    message: Optional[str] = None
    output_file: Optional[str] = None


@router.post("/generate", response_model=JobResponse)
async def generate_presentation(
    file: UploadFile = File(...),
    output_type: str = Form(default="slides"),  # slides or poster
    template: str = Form(default="academic"),
    slide_count: int = Form(default=10),
    additional_prompt: str = Form(default=""),
):
    """
    Paper2Slides /generate 엔드포인트를 호출하여 프레젠테이션 생성을 시작합니다.
    """

    # Map slide_count to length
    length_map = {5: "short", 10: "medium", 15: "long", 20: "long"}
    length = length_map.get(slide_count, "medium")

    # Map template to style (Paper2Slides 지원: academic, doraemon, or custom)
    style_map = {
        "asan-medical": "academic",
        "academic": "academic",
        "research": "academic",
        "clinical": "academic",
        "minimal": "academic",
    }
    style = style_map.get(template, "academic")

    try:
        file_content = await file.read()

        async with httpx.AsyncClient(timeout=30.0) as client:
            files = {"file": (file.filename, file_content, file.content_type or "application/pdf")}
            data = {
                "output_type": output_type,
                "style": style,
                "length": length,
            }

            response = await client.post(
                f"{PAPER2SLIDES_API_URL}/generate",
                files=files,
                data=data,
            )
            response.raise_for_status()
            result = response.json()

            return JobResponse(
                job_id=result["job_id"],
                status=result["status"],
                message=result.get("message", "Paper2Slides 작업이 시작되었습니다."),
            )

    except httpx.TimeoutException:
        raise HTTPException(status_code=504, detail="Paper2Slides 서버 응답 시간 초과")
    except httpx.ConnectError:
        raise HTTPException(
            status_code=503,
            detail="Paper2Slides 서버에 연결할 수 없습니다. SSH 터널(localhost:29001)을 확인해주세요.",
        )
    except httpx.HTTPStatusError as e:
        detail = str(e)
        try:
            detail = e.response.json().get("detail", detail)
        except Exception:
            pass
        raise HTTPException(status_code=e.response.status_code, detail=detail)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Paper2Slides 호출 실패: {str(e)}")


@router.get("/status/{job_id}", response_model=JobStatus)
async def get_job_status(job_id: str):
    """Paper2Slides /status 프록시"""
    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            response = await client.get(f"{PAPER2SLIDES_API_URL}/status/{job_id}")
            response.raise_for_status()
            return JobStatus(**response.json())
    except httpx.HTTPStatusError as e:
        if e.response.status_code == 404:
            raise HTTPException(status_code=404, detail="작업을 찾을 수 없습니다")
        raise HTTPException(status_code=e.response.status_code, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/download/{job_id}")
async def download_output(job_id: str):
    """Paper2Slides /download 프록시 - 생성된 PDF 다운로드"""
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            response = await client.get(
                f"{PAPER2SLIDES_API_URL}/download/{job_id}",
                follow_redirects=True,
            )
            response.raise_for_status()

            content_type = response.headers.get("content-type", "application/pdf")

            return StreamingResponse(
                iter([response.content]),
                media_type=content_type,
                headers={
                    "Content-Disposition": f'attachment; filename="presentation_{job_id}.pdf"'
                },
            )
    except httpx.HTTPStatusError as e:
        if e.response.status_code == 404:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다")
        raise HTTPException(status_code=e.response.status_code, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/pptx/{job_id}")
async def download_pptx(job_id: str):
    """생성된 PDF를 페이지별 이미지로 변환하여 PPTX로 다운로드"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="PyMuPDF가 설치되지 않았습니다. pip install pymupdf",
        )

    try:
        # 1. Paper2Slides에서 PDF 다운로드
        async with httpx.AsyncClient(timeout=120.0) as client:
            response = await client.get(
                f"{PAPER2SLIDES_API_URL}/download/{job_id}",
                follow_redirects=True,
            )
            response.raise_for_status()
            pdf_bytes = response.content

        # 2. PDF를 페이지별 이미지로 변환 (PyMuPDF)
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        page_images: list[bytes] = []
        for page in doc:
            # 고해상도 렌더링 (2x zoom = 144 DPI)
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            page_images.append(pix.tobytes("png"))
        doc.close()

        if not page_images:
            raise HTTPException(status_code=500, detail="PDF에서 이미지를 추출할 수 없습니다")

        # 3. python-pptx로 PPTX 생성 (16:9 와이드스크린)
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]  # Blank layout

        for img_data in page_images:
            slide = prs.slides.add_slide(blank_layout)
            img_stream = io.BytesIO(img_data)
            slide.shapes.add_picture(
                img_stream,
                Emu(0), Emu(0),
                prs.slide_width, prs.slide_height,
            )

        # 4. PPTX를 메모리에 저장 후 반환
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        filename = f"presentation_{job_id}.pptx"
        return StreamingResponse(
            pptx_buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"'
            },
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX 변환 실패: {str(e)}")


@router.get("/health")
async def health_check():
    """Paper2Slides 연결 상태 확인"""
    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            response = await client.get(f"{PAPER2SLIDES_API_URL}/health")
            if response.status_code == 200:
                return {
                    "status": "healthy",
                    "backend": "Paper2Slides (HKUDS)",
                    "llm": "Qwen3-32B-AWQ",
                    "endpoint": PAPER2SLIDES_API_URL,
                }
    except Exception:
        pass
    return {
        "status": "unhealthy",
        "backend": "Paper2Slides (HKUDS)",
        "endpoint": PAPER2SLIDES_API_URL,
    }
